#!/usr/bin/env python
"""
Analyze images in PowerPoint slides to understand what content we're missing
"""

import pptx
from pathlib import Path
from PIL import Image
import io

def analyze_slide_images(pptx_path, slide_num):
    """Extract and analyze images from a specific slide"""
    prs = pptx.Presentation(str(pptx_path))
    
    if len(prs.slides) < slide_num:
        print(f"File only has {len(prs.slides)} slides")
        return
    
    slide = prs.slides[slide_num - 1]  # 0-indexed
    
    print(f"Analyzing Slide {slide_num}")
    print("=" * 60)
    
    # Count different shape types
    text_shapes = 0
    image_shapes = 0
    other_shapes = 0
    
    for shape_idx, shape in enumerate(slide.shapes):
        # Text content
        if hasattr(shape, 'text') and shape.text.strip():
            text_shapes += 1
            print(f"Text Shape {shape_idx}: {shape.text.strip()[:100]}")
        
        # Images
        if shape.shape_type == 13:  # Picture type
            image_shapes += 1
            print(f"\nImage Shape {shape_idx}:")
            print(f"  - Position: ({shape.left}, {shape.top})")
            print(f"  - Size: {shape.width} x {shape.height}")
            
            # Try to extract image
            if hasattr(shape, 'image'):
                image = shape.image
                print(f"  - Image format: {image.ext}")
                print(f"  - Image size: {len(image.blob)} bytes")
                
                # Save image for manual inspection
                image_filename = f"slide_{slide_num}_image_{shape_idx}.{image.ext}"
                with open(image_filename, 'wb') as f:
                    f.write(image.blob)
                print(f"  - Saved as: {image_filename}")
                
                # Try to open with PIL to get dimensions
                try:
                    pil_image = Image.open(io.BytesIO(image.blob))
                    print(f"  - Image dimensions: {pil_image.size}")
                    print(f"  - Image mode: {pil_image.mode}")
                except Exception as e:
                    print(f"  - Could not analyze image: {e}")
        else:
            other_shapes += 1
    
    print(f"\nSummary:")
    print(f"  - Text shapes: {text_shapes}")
    print(f"  - Image shapes: {image_shapes}")
    print(f"  - Other shapes: {other_shapes}")
    print(f"  - Total shapes: {len(slide.shapes)}")
    
    return image_shapes > 0

if __name__ == "__main__":
    pptx_path = Path("C:/IU/Level 4/Discrete Maths/Lecture Notes/CO3.pptx")
    
    # Analyze slide 37
    print("\n" + "=" * 60)
    has_images = analyze_slide_images(pptx_path, 37)
    
    # Also check a few other slides
    print("\n" + "=" * 60)
    print("\nChecking other slides for comparison:")
    for slide_num in [34, 35, 36]:
        print(f"\nSlide {slide_num}:")
        slide = pptx.Presentation(str(pptx_path)).slides[slide_num - 1]
        image_count = sum(1 for shape in slide.shapes if shape.shape_type == 13)
        text_count = sum(1 for shape in slide.shapes if hasattr(shape, 'text') and shape.text.strip())
        print(f"  - {text_count} text shapes, {image_count} image shapes")