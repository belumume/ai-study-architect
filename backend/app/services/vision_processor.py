"""
Vision-based content processing using Claude and OpenAI APIs
Processes images, extracts text, and understands visual content
"""

import os
import base64
import logging
from typing import Dict, Any, Optional, List
from pathlib import Path
import io
from PIL import Image
from app.services.claude_service import claude_service

# OpenAI service will be imported only if needed
try:
    from app.services.openai_service import openai_service
    HAS_OPENAI = True
except ImportError:
    HAS_OPENAI = False
    openai_service = None

logger = logging.getLogger(__name__)


class VisionProcessor:
    """Process visual content using AI vision APIs"""
    
    def __init__(self):
        """Initialize vision processor"""
        self.max_image_size = (2048, 2048)  # Max dimensions for API
        self.supported_formats = ['png', 'jpg', 'jpeg', 'gif', 'webp']
        
    def prepare_image(self, image_data: bytes) -> tuple[bytes, str]:
        """
        Prepare image for API submission
        
        Args:
            image_data: Raw image bytes
            
        Returns:
            Tuple of (processed_bytes, mime_type)
        """
        try:
            # Open image
            img = Image.open(io.BytesIO(image_data))
            
            # Convert RGBA to RGB if needed
            if img.mode in ('RGBA', 'LA', 'P'):
                rgb_img = Image.new('RGB', img.size, (255, 255, 255))
                rgb_img.paste(img, mask=img.split()[-1] if img.mode == 'RGBA' else None)
                img = rgb_img
            
            # Resize if too large
            if img.width > self.max_image_size[0] or img.height > self.max_image_size[1]:
                img.thumbnail(self.max_image_size, Image.Resampling.LANCZOS)
            
            # Save to bytes
            output = io.BytesIO()
            img.save(output, format='PNG')
            output.seek(0)
            
            return output.read(), 'image/png'
            
        except Exception as e:
            logger.error(f"Error preparing image: {e}")
            return image_data, 'image/png'
    
    async def extract_with_claude(self, image_data: bytes, prompt: Optional[str] = None) -> Dict[str, Any]:
        """
        Extract content using Claude Vision
        
        Args:
            image_data: Image bytes
            prompt: Optional specific extraction prompt
            
        Returns:
            Extracted content and metadata
        """
        try:
            # Prepare image
            processed_image, mime_type = self.prepare_image(image_data)
            
            # Encode to base64
            image_base64 = base64.b64encode(processed_image).decode('utf-8')
            
            # Build prompt
            if not prompt:
                prompt = """Please analyze this image and:
1. Extract ALL text content you can see
2. If it contains mathematical formulas, write them in LaTeX format
3. If it contains a table, format it clearly
4. If it contains a diagram, describe its structure
5. Identify the type of content (text, formula, table, diagram, mixed)

Provide the response in this format:
EXTRACTED TEXT:
[all text content]

CONTENT TYPE: [text/formula/table/diagram/mixed]

DESCRIPTION:
[brief description of what the image contains]"""
            
            # Call Claude Vision API
            messages = [
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "image",
                            "source": {
                                "type": "base64",
                                "media_type": mime_type,
                                "data": image_base64
                            }
                        },
                        {
                            "type": "text",
                            "text": prompt
                        }
                    ]
                }
            ]
            
            response = await claude_service.create_completion(
                messages=messages,
                max_tokens=2000,
                system="You are a vision AI extracting content from educational materials."
            )
            
            # Parse response
            content = response.get('content', [{}])[0].get('text', '')
            
            # Extract sections
            extracted_text = ""
            content_type = "unknown"
            description = ""
            
            if "EXTRACTED TEXT:" in content:
                parts = content.split("EXTRACTED TEXT:")[1]
                if "CONTENT TYPE:" in parts:
                    extracted_text = parts.split("CONTENT TYPE:")[0].strip()
                    remaining = parts.split("CONTENT TYPE:")[1]
                    if "DESCRIPTION:" in remaining:
                        content_type = remaining.split("DESCRIPTION:")[0].strip()
                        description = remaining.split("DESCRIPTION:")[1].strip()
            else:
                # Fallback: treat entire response as extracted text
                extracted_text = content
            
            return {
                "success": True,
                "text": extracted_text,
                "content_type": content_type,
                "description": description,
                "processor": "claude_vision"
            }
            
        except Exception as e:
            logger.error(f"Claude Vision extraction failed: {e}")
            return {
                "success": False,
                "text": "",
                "error": str(e),
                "processor": "claude_vision"
            }
    
    async def extract_with_openai(self, image_data: bytes, prompt: Optional[str] = None) -> Dict[str, Any]:
        """
        Extract content using OpenAI Vision (GPT-4V)
        
        Args:
            image_data: Image bytes
            prompt: Optional specific extraction prompt
            
        Returns:
            Extracted content and metadata
        """
        try:
            # Prepare image
            processed_image, _ = self.prepare_image(image_data)
            
            # Encode to base64
            image_base64 = base64.b64encode(processed_image).decode('utf-8')
            
            # Build prompt
            if not prompt:
                prompt = """Analyze this image and extract:
1. ALL visible text
2. Mathematical formulas (in LaTeX)
3. Tables (formatted clearly)
4. Diagram descriptions
Return: extracted text, content type, and brief description."""
            
            # Call OpenAI Vision API
            messages = [
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "text",
                            "text": prompt
                        },
                        {
                            "type": "image_url",
                            "image_url": {
                                "url": f"data:image/png;base64,{image_base64}"
                            }
                        }
                    ]
                }
            ]
            
            response = await openai_service.create_completion(
                messages=messages,
                model="gpt-4-vision-preview",
                max_tokens=2000
            )
            
            content = response.choices[0].message.content
            
            return {
                "success": True,
                "text": content,
                "content_type": "extracted",
                "description": "OpenAI Vision extraction",
                "processor": "openai_vision"
            }
            
        except Exception as e:
            logger.error(f"OpenAI Vision extraction failed: {e}")
            return {
                "success": False,
                "text": "",
                "error": str(e),
                "processor": "openai_vision"
            }
    
    async def extract_from_image(self, image_data: bytes, use_fallback: bool = True) -> Dict[str, Any]:
        """
        Extract content from image using available vision APIs
        
        Args:
            image_data: Image bytes
            use_fallback: Try OpenAI if Claude fails
            
        Returns:
            Extracted content and metadata
        """
        # Try Claude first (better for educational content)
        anthropic_key = os.getenv("ANTHROPIC_API_KEY")
        if anthropic_key:
            result = await self.extract_with_claude(image_data)
            if result["success"]:
                return result
        
        # Fallback to OpenAI
        openai_key = os.getenv("OPENAI_API_KEY")
        if use_fallback and HAS_OPENAI and openai_key:
            result = await self.extract_with_openai(image_data)
            if result["success"]:
                return result
        
        # No vision API available or all failed
        return {
            "success": False,
            "text": "Vision-based text extraction not available. Please ensure AI API keys are configured.",
            "content_type": "error",
            "description": "No vision API available",
            "processor": "none"
        }
    
    def extract_from_image_sync(self, image_data: bytes, use_fallback: bool = True) -> Dict[str, Any]:
        """
        Synchronous wrapper for extract_from_image
        Safe to call from sync contexts like content_processor
        
        Args:
            image_data: Image bytes
            use_fallback: Try OpenAI if Claude fails
            
        Returns:
            Extracted content and metadata
        """
        import asyncio
        import nest_asyncio
        
        # Allow nested event loops (for FastAPI compatibility)
        nest_asyncio.apply()
        
        try:
            # Try to get existing event loop
            try:
                loop = asyncio.get_running_loop()
                # We're in an async context, create a new loop in thread
                import concurrent.futures
                with concurrent.futures.ThreadPoolExecutor() as executor:
                    future = executor.submit(asyncio.run, self.extract_from_image(image_data, use_fallback))
                    return future.result()
            except RuntimeError:
                # No event loop running, we can run directly
                return asyncio.run(self.extract_from_image(image_data, use_fallback))
        except Exception as e:
            logger.error(f"Sync extraction failed: {e}")
            return {
                "success": False,
                "text": f"Vision extraction error: {str(e)}",
                "content_type": "error",
                "description": "Extraction failed",
                "processor": "none"
            }
    
    async def extract_from_pptx_images(self, pptx_path: Path) -> List[Dict[str, Any]]:
        """
        Extract content from all images in a PowerPoint
        
        Args:
            pptx_path: Path to PowerPoint file
            
        Returns:
            List of extraction results per slide
        """
        import pptx
        
        results = []
        try:
            prs = pptx.Presentation(str(pptx_path))
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_images = []
                
                for shape in slide.shapes:
                    if shape.shape_type == 13:  # Picture type
                        if hasattr(shape, 'image'):
                            image_data = shape.image.blob
                            
                            # Extract content from image
                            extraction = await self.extract_from_image(image_data)
                            extraction["slide_number"] = slide_num
                            extraction["shape_position"] = (shape.left, shape.top)
                            
                            slide_images.append(extraction)
                
                if slide_images:
                    results.append({
                        "slide_number": slide_num,
                        "images": slide_images
                    })
            
        except Exception as e:
            logger.error(f"Error extracting from PowerPoint images: {e}")
        
        return results
    
    async def extract_from_pdf_images(self, pdf_path: Path) -> List[Dict[str, Any]]:
        """
        Extract content from images in PDF
        
        Args:
            pdf_path: Path to PDF file
            
        Returns:
            List of extraction results per page
        """
        import fitz  # PyMuPDF
        
        results = []
        try:
            pdf = fitz.open(str(pdf_path))
            
            for page_num, page in enumerate(pdf, 1):
                page_images = []
                
                # Get images from page
                image_list = page.get_images()
                
                for img_index, img in enumerate(image_list):
                    # Extract image
                    xref = img[0]
                    pix = fitz.Pixmap(pdf, xref)
                    
                    if pix.n - pix.alpha < 4:  # GRAY or RGB
                        img_data = pix.tobytes("png")
                    else:  # CMYK
                        pix = fitz.Pixmap(fitz.csRGB, pix)
                        img_data = pix.tobytes("png")
                    
                    # Extract content from image
                    extraction = await self.extract_from_image(img_data)
                    extraction["page_number"] = page_num
                    extraction["image_index"] = img_index
                    
                    page_images.append(extraction)
                    pix = None
                
                if page_images:
                    results.append({
                        "page_number": page_num,
                        "images": page_images
                    })
            
            pdf.close()
            
        except Exception as e:
            logger.error(f"Error extracting from PDF images: {e}")
        
        return results


# Global instance
vision_processor = VisionProcessor()