"""
Content processing service for extracting text from various file formats
"""

import os
import logging
import mimetypes
from pathlib import Path
from typing import Optional, Dict, Any, List
import PyPDF2
import docx
import pptx
import magic
from PIL import Image

# Optional OCR support
try:
    import pytesseract
    HAS_OCR = True
except ImportError:
    HAS_OCR = False

from app.core.config import settings
from app.core.exceptions import ContentProcessingError

# Import vision processor for image extraction
try:
    from app.services.vision_processor import vision_processor
    HAS_VISION = True
except ImportError:
    HAS_VISION = False

logger = logging.getLogger(__name__)


class ContentProcessor:
    """Service for processing uploaded content files"""
    
    SUPPORTED_FORMATS = {
        'application/pdf': 'pdf',
        'application/vnd.openxmlformats-officedocument.wordprocessingml.document': 'docx',
        'application/vnd.openxmlformats-officedocument.presentationml.presentation': 'pptx',
        'text/plain': 'txt',
        'text/markdown': 'md',
        'image/jpeg': 'image',
        'image/png': 'image',
        'image/jpg': 'image'
    }
    
    def __init__(self):
        """Initialize content processor"""
        self.upload_dir = Path(settings.UPLOAD_DIR)
        
    def process_file(self, file_path: str, content_type: str) -> Dict[str, Any]:
        """
        Process a file and extract its content
        
        Args:
            file_path: Path to the file
            content_type: MIME type of the file
            
        Returns:
            Dict containing extracted text and metadata
        """
        try:
            file_path = Path(file_path)
            if not file_path.exists():
                raise ContentProcessingError(f"File not found: {file_path}")
            
            # Detect actual file type using magic
            actual_mime = magic.from_file(str(file_path), mime=True)
            logger.info(f"Processing file: {file_path.name}, declared type: {content_type}, actual type: {actual_mime}")
            
            # Extract text based on file type
            text = ""
            metadata = {
                "file_size": file_path.stat().st_size,
                "detected_type": actual_mime,
                "declared_type": content_type
            }
            
            # Handle PowerPoint files detected as zip
            if actual_mime == "application/zip" and file_path.suffix.lower() == ".pptx":
                actual_mime = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
            
            if actual_mime == 'application/pdf':
                text, pdf_metadata = self._extract_pdf_text(file_path)
                metadata.update(pdf_metadata)
            elif actual_mime in ['application/vnd.openxmlformats-officedocument.wordprocessingml.document']:
                text = self._extract_docx_text(file_path)
            elif actual_mime in ['application/vnd.openxmlformats-officedocument.presentationml.presentation']:
                text = self._extract_pptx_text(file_path)
            elif actual_mime.startswith('text/'):
                text = self._extract_text_file(file_path)
            elif actual_mime.startswith('image/'):
                text = self._extract_image_text(file_path)
                metadata["extraction_method"] = "OCR"
            else:
                logger.warning(f"Unsupported file type: {actual_mime}")
                text = f"Unable to extract text from {actual_mime} file"
                
            # Calculate text statistics
            metadata["text_length"] = len(text)
            metadata["word_count"] = len(text.split())
            metadata["line_count"] = len(text.split('\n'))
            
            return {
                "text": text,
                "metadata": metadata,
                "success": True
            }
            
        except Exception as e:
            logger.error(f"Error processing file {file_path}: {str(e)}", exc_info=True)
            return {
                "text": "",
                "metadata": {"error": str(e)},
                "success": False
            }
    
    def _extract_pdf_text(self, file_path: Path) -> tuple[str, Dict[str, Any]]:
        """Extract text from PDF file including OCR for scanned pages"""
        text_parts = []
        metadata = {"page_count": 0, "has_images": False}
        
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                metadata["page_count"] = len(pdf_reader.pages)
                
                for page_num, page in enumerate(pdf_reader.pages):
                    page_content = [f"--- Page {page_num + 1} ---"]
                    
                    # Extract regular text
                    try:
                        page_text = page.extract_text()
                        if page_text.strip():
                            page_content.append(page_text)
                    except Exception as e:
                        logger.warning(f"Error extracting text from page {page_num + 1}: {e}")
                    
                    # Check if page might be scanned (no extractable text)
                    if len(page_content) == 1:  # Only header, no text
                        # Try to extract as image if vision AI is available
                        if HAS_VISION:
                            try:
                                # Note: This is a simplified approach
                                # Full implementation would use PyMuPDF or pdf2image
                                page_content.append("[Page appears to be scanned - vision extraction needed]")
                                metadata["has_images"] = True
                            except Exception as e:
                                logger.warning(f"Could not check for images on page {page_num + 1}: {e}")
                    
                    text_parts.append('\n'.join(page_content))
                
                # Extract metadata
                if pdf_reader.metadata:
                    metadata["title"] = pdf_reader.metadata.get('/Title', '')
                    metadata["author"] = pdf_reader.metadata.get('/Author', '')
                    metadata["subject"] = pdf_reader.metadata.get('/Subject', '')
                    
        except Exception as e:
            logger.error(f"Error reading PDF: {e}")
            raise ContentProcessingError(f"Failed to process PDF: {str(e)}")
        
        return '\n\n'.join(text_parts), metadata
    
    def _extract_docx_text(self, file_path: Path) -> str:
        """Extract text from DOCX file"""
        try:
            doc = docx.Document(file_path)
            text_parts = []
            
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text_parts.append(paragraph.text)
            
            # Also extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text.strip())
                    if row_text:
                        text_parts.append(' | '.join(row_text))
            
            return '\n\n'.join(text_parts)
            
        except Exception as e:
            logger.error(f"Error reading DOCX: {e}")
            raise ContentProcessingError(f"Failed to process DOCX: {str(e)}")
    
    def _extract_pptx_text(self, file_path: Path) -> str:
        """Extract text from PPTX file including images with vision AI"""
        try:
            prs = pptx.Presentation(file_path)
            text_parts = []
            
            for slide_num, slide in enumerate(prs.slides):
                slide_text = [f"--- Slide {slide_num + 1} ---"]
                image_count = 0
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)
                    
                    # Extract text from tables
                    if hasattr(shape, 'has_table') and shape.has_table:
                        for row in shape.table.rows:
                            row_text = []
                            for cell in row.cells:
                                if cell.text.strip():
                                    row_text.append(cell.text.strip())
                            if row_text:
                                slide_text.append(' | '.join(row_text))
                    
                    # NEW: Extract content from images using vision AI
                    if shape.shape_type == 13 and HAS_VISION:  # Picture type
                        image_count += 1
                        try:
                            if hasattr(shape, 'image'):
                                image_data = shape.image.blob
                                # Use vision API to extract content
                                import asyncio
                                loop = asyncio.new_event_loop()
                                asyncio.set_event_loop(loop)
                                extraction = loop.run_until_complete(
                                    vision_processor.extract_from_image(image_data)
                                )
                                loop.close()
                                
                                if extraction['success'] and extraction['text']:
                                    slide_text.append(f"\n[Image {image_count} content:]")
                                    slide_text.append(extraction['text'])
                                    if extraction.get('description'):
                                        slide_text.append(f"[Description: {extraction['description']}]")
                        except Exception as e:
                            logger.warning(f"Could not extract image from slide {slide_num + 1}: {e}")
                            slide_text.append(f"[Image {image_count}: Unable to extract content]")
                
                if len(slide_text) > 1:  # More than just the slide header
                    text_parts.append('\n'.join(slide_text))
            
            return '\n\n'.join(text_parts)
            
        except Exception as e:
            logger.error(f"Error reading PPTX: {e}")
            raise ContentProcessingError(f"Failed to process PPTX: {str(e)}")
    
    def _extract_text_file(self, file_path: Path) -> str:
        """Extract text from plain text file"""
        try:
            # Try different encodings
            encodings = ['utf-8', 'latin-1', 'cp1252']
            
            for encoding in encodings:
                try:
                    with open(file_path, 'r', encoding=encoding) as file:
                        return file.read()
                except UnicodeDecodeError:
                    continue
            
            # If all encodings fail, read as binary and decode with errors ignored
            with open(file_path, 'rb') as file:
                return file.read().decode('utf-8', errors='ignore')
                
        except Exception as e:
            logger.error(f"Error reading text file: {e}")
            raise ContentProcessingError(f"Failed to process text file: {str(e)}")
    
    def _extract_image_text(self, file_path: Path) -> str:
        """Extract text from image using vision AI or OCR"""
        # Try vision AI first (much better than local OCR)
        if HAS_VISION:
            try:
                with open(file_path, 'rb') as f:
                    image_data = f.read()
                
                # Use vision API for extraction
                import asyncio
                loop = asyncio.new_event_loop()
                asyncio.set_event_loop(loop)
                extraction = loop.run_until_complete(
                    vision_processor.extract_from_image(image_data)
                )
                loop.close()
                
                if extraction['success'] and extraction['text']:
                    result = extraction['text']
                    if extraction.get('description'):
                        result += f"\n\n[Image Description: {extraction['description']}]"
                    return result
            except Exception as e:
                logger.error(f"Vision AI extraction failed: {e}")
        
        # Fallback to local OCR if available
        if HAS_OCR:
            try:
                # Open image with PIL
                image = Image.open(file_path)
                
                # Convert to RGB if necessary
                if image.mode != 'RGB':
                    image = image.convert('RGB')
                
                # Extract text using OCR
                text = pytesseract.image_to_string(image)
                
                if not text.strip():
                    return "No text found in image"
                
                return text
                
            except Exception as e:
                logger.error(f"Error performing OCR on image: {e}")
        
        # Neither vision AI nor OCR available
        return "Text extraction from images requires AI vision API keys to be configured."
    
    def extract_chunks(self, text: str, chunk_size: int = 1000, overlap: int = 200) -> List[str]:
        """
        Split text into overlapping chunks for processing
        
        Args:
            text: The text to chunk
            chunk_size: Size of each chunk in characters
            overlap: Number of overlapping characters between chunks
            
        Returns:
            List of text chunks
        """
        if not text:
            return []
        
        chunks = []
        start = 0
        text_length = len(text)
        
        while start < text_length:
            # Find the end of the chunk
            end = start + chunk_size
            
            # If this is not the last chunk, try to break at a sentence or word boundary
            if end < text_length:
                # Look for sentence boundaries
                for delimiter in ['. ', '.\n', '! ', '!\n', '? ', '?\n']:
                    delimiter_pos = text.rfind(delimiter, start, end)
                    if delimiter_pos != -1:
                        end = delimiter_pos + len(delimiter)
                        break
                else:
                    # No sentence boundary found, try word boundary
                    space_pos = text.rfind(' ', start, end)
                    if space_pos != -1:
                        end = space_pos + 1
            
            # Extract chunk
            chunk = text[start:end].strip()
            if chunk:
                chunks.append(chunk)
            
            # Move start position
            start = end - overlap if end < text_length else text_length
        
        return chunks


# Global instance
content_processor = ContentProcessor()